"""Fill the Review-2 PPT template in place.

Same rule as the report: the template is the design. Every bullet reuses the
formatting of the paragraph already in its placeholder, and no layout, position
or font is changed. Slide 2 is left as-is because it needs a screenshot of the
guide's approval mail, which only the student can supply.

    python tools/build_ppt.py
"""
from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
SRC = Path(r"C:\Users\M V S Akhil Teja\Downloads\4 Review 2 PPT Template.pptx")
OUT = ROOT / "report" / "BCSE497J_Review2_23BCB0135.pptx"
FIG = ROOT / "report" / "figures"

TITLE = ("Deep Learning-Based Classification of Standard and "
         "Non-Standard Intrapartum Ultrasound Images")

SLIDES = {
    3: [  # Aim
        (0, "To build a deep learning classifier that decides whether an "
            "intrapartum transperineal ultrasound frame is a standard plane, "
            "i.e. whether the pubic symphysis and fetal head are resolved "
            "well enough to measure the angle of progression."),
        (0, "To use that classifier to settle a methodological question the "
            "literature leaves open: does frame sampling density itself "
            "improve performance, or is the reported advantage of video data "
            "simply the effect of having more training samples?"),
    ],
    4: [  # Abstract
        (0, "Standard-plane selection during labour is manual, "
            "operator-dependent and performed under time pressure."),
        (0, "Prior work reports that dense video frames beat curated images, "
            "but changes the corpus and the sample count together, so the "
            "cause is unidentifiable."),
        (0, "Both arms are built from one corpus (IUGC 2024, 774 videos, "
            "65,531 frames), varying only frames drawn per video, so probe, "
            "anatomy, scanner and annotator are constant by construction."),
        (0, "ResNet-18 trained under 13 sampling conditions; ResNet-18 + "
            "BiLSTM for the temporal arm; all scored on the same 8,665-frame "
            "official test set."),
        (0, "Result: no density effect. The best single model used 434 "
            "training frames. Across three seeds the matched-budget "
            "comparison reverses sign, so seed variance exceeds the effect."),
    ],
    5: [  # Literature Review
        (0, "Chen et al. (2015) - transferred CNN features localise fetal "
            "standard planes far better than hand-crafted descriptors."),
        (0, "Baumgartner et al. (2017), SonoNet - real-time detection of 13 "
            "fetal standard planes from image-level labels; reference "
            "architecture for the task."),
        (0, "Burgos-Artizzu et al. (2020) - FETAL_PLANES_DB, ~12,400 curated "
            "images; the canonical curated still-image regime."),
        (0, "Ghi et al. (2018), ISUOG guidelines; Kalache et al. (2009) - "
            "define the transperineal views and the angle of progression, "
            "i.e. what the positive class clinically means."),
        (0, "Two patterns recur: image-vs-video comparisons change corpus and "
            "sample count together, and temporal models are benchmarked "
            "against unsmoothed frame-wise baselines."),
    ],
    6: [  # Research Gap
        (0, "Gap 1 - Sampling density is never isolated from sample count. No "
            "prior work reports a budget-matched comparison where both arms "
            "receive the same number of training frames."),
        (0, "Gap 2 - Temporal models are compared against unsmoothed "
            "baselines. Standard planes occur in contiguous runs, so a "
            "zero-parameter moving average already recovers much of that "
            "structure."),
        (0, "Gap 3 - Run-to-run variance is not reported. On corpora of a few "
            "hundred patients, results come from a single training run; if "
            "seed variance is comparable to the effect, the comparison cannot "
            "support its conclusion."),
    ],
    7: [  # Objectives - kept to five short lines: the SDG/outcomes strip
          # and its icons start at 4.83in, so the body has ~2.8in of room.
        (0, "Build a verified frame-level dataset from IUGC 2024 (label join "
            "rate 1.0, integrity checks as assertions)."),
        (0, "Train and evaluate a ResNet-18 standard-plane classifier on the "
            "official held-out test split."),
        (0, "Quantify sampling density over 13 conditions, and separate "
            "density from sample count with budget- and prior-matched arms."),
        (0, "Test temporal modelling against a smoothed frame-wise baseline, "
            "not an unsmoothed one."),
        (0, "Measure seed-to-seed variance; deploy the classifier as a "
            "verified browser demonstrator."),
    ],
    9: [  # Functional Requirements
        (0, "Frame extraction - decode each video once to pre-resized frames "
            "through a Unicode-safe write path."),
        (0, "Label parsing - handle per-video index lists and the ALL / NONE "
            "whole-video sentinels; fail loudly, never silently."),
        (0, "Index construction - one frame-level index; refuse to write it "
            "if the label join rate is below 1.0."),
        (0, "Video-level splitting - partition by video, never by frame; "
            "assert no video appears in two splits."),
        (0, "Sampling conditions - build a training set for any k per video, "
            "stride, frame budget or target class prior."),
        (0, "Training and evaluation - one trainer for both arms; report "
            "macro-F1, balanced accuracy, AUPRC and MCC with video-level "
            "bootstrap intervals."),
        (0, "Inference - return the standard-plane probability and the "
            "decision at the frozen threshold for a submitted frame."),
    ],
    10: [  # Modules
        (0, "Data preparation - extract_frames, build_index: decoding, "
            "sentinel-aware label parsing, integrity assertions."),
        (0, "Splitting - splits: video-level official split, leakage "
            "assertions, label-transition diagnostics."),
        (0, "Sampling - sampling: k-per-video, stride, budget matching, prior "
            "matching. This module is the study's independent variable."),
        (0, "Datasets - datasets: frame and clip views, ultrasound-specific "
            "augmentation, transition splicing for the temporal arm."),
        (0, "Models - models: one shared ResNet-18 encoder, two heads "
            "(frame-wise and BiLSTM)."),
        (0, "Training and evaluation - engine, metrics, smoothing: shared "
            "trainer, imbalance-aware metrics, post-hoc smoothing baseline."),
        (0, "Deployment - export_onnx and a static page running "
            "onnxruntime-web for local, in-browser inference."),
    ],
    11: [  # Experiments and Results
        (0, "Arm 1, 13 conditions, seed 0: sparse arm 0.588-0.687 macro-F1 "
            "over 434 to 8,638 frames; dense arm 0.506-0.637 over 6,840 to "
            "53,996 frames. No trend in either. Best result: 434 frames."),
        (0, "Arm 2, same dense condition: frame-wise 0.5785 raw, 0.5453 "
            "smoothed; BiLSTM 0.6695 with splicing, 0.6496 without. The "
            "temporal arm wins, but the ablation shows splicing is not why."),
        (0, "Three seeds: budget-matched delta -0.106 / +0.013 / -0.093; "
            "prior-matched -0.168 / +0.121 / -0.045. Both reverse sign; both "
            "standard deviations exceed their means."),
        (0, "One condition varies by 0.231 macro-F1 across seeds on identical "
            "data with an identical recipe."),
        (0, "Post-hoc smoothing hurt in 12 of 13 conditions: it is tuned on "
            "validation, where runs are long (median 33 frames), and applied "
            "to test, where they are short (median 12)."),
    ],
    12: [  # Conclusion
        (0, "Frame sampling density produced no measurable benefit on this "
            "corpus. The strongest single model was trained on 434 frames, "
            "one per video, against 53,996 for the dense arm."),
        (0, "The matched-budget result does not replicate across seeds, so "
            "Protocol B is unresolved rather than answered."),
        (0, "The principal finding is a measurement one: at 434 training "
            "videos with this train-test shift, seed variance reaches about "
            "+/-0.12 macro-F1, exceeding any density effect. Single-run "
            "comparisons at this scale cannot support conclusions about "
            "density - and that is what prior work reports."),
        (0, "Remaining work: Grad-CAM anatomical attention check, "
            "leave-one-centre-out split, paired video bootstrap, and AoP/HSD "
            "downstream error using the corpus landmark annotations."),
    ],
    13: [  # References
        (0, "[1] C. F. Baumgartner et al., \"SonoNet: Real-time detection and "
            "localisation of fetal standard scan planes in freehand "
            "ultrasound,\" IEEE Trans. Med. Imaging, vol. 36, no. 11, "
            "pp. 2204-2215, 2017."),
        (0, "[2] H. Chen et al., \"Standard plane localization in fetal "
            "ultrasound via domain transferred deep neural networks,\" IEEE "
            "J. Biomed. Health Inform., vol. 19, no. 5, pp. 1627-1636, 2015."),
        (0, "[3] X. P. Burgos-Artizzu et al., \"Evaluation of deep "
            "convolutional neural networks for automatic classification of "
            "common maternal fetal ultrasound planes,\" Sci. Rep., vol. 10, "
            "p. 10200, 2020."),
        (0, "[4] T. Ghi et al., \"ISUOG Practice Guidelines: intrapartum "
            "ultrasound,\" Ultrasound Obstet. Gynecol., vol. 52, no. 1, "
            "pp. 128-139, 2018."),
        (0, "[5] K. He, X. Zhang, S. Ren and J. Sun, \"Deep residual learning "
            "for image recognition,\" in Proc. IEEE CVPR, 2016, "
            "pp. 770-778."),
        (0, "[6] IUGC 2024 dataset, Zenodo, doi: 10.5281/zenodo.17655183, "
            "CC-BY-4.0."),
    ],
}

SDG_LINE = ("Identified:   Sustainable Development Goal 3 - Good Health and "
            "Well-being                    Outcomes: Conference (Scopus) / "
            "Product (open-source classifier and browser demonstrator)")


def set_bullets(tf, items):
    """Replace a text frame's paragraphs, reusing the first one's formatting."""
    p0 = tf.paragraphs[0]
    proto = copy.deepcopy(p0._p)
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)

    def fill(par, level, text):
        if par.runs:
            par.runs[0].text = text
            for r in list(par.runs)[1:]:
                r._r.getparent().remove(r._r)
        else:
            par.add_run().text = text
        par.level = level

    fill(tf.paragraphs[0], *items[0])
    prev = tf.paragraphs[0]._p
    for level, text in items[1:]:
        new = copy.deepcopy(proto)
        prev.addnext(new)
        prev = new
        # tf.paragraphs re-reads the XML, so the one just inserted is last
        fill(tf.paragraphs[-1], level, text)


def content_ph(slide):
    """The body text placeholder on a 'Title and Content' slide."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if sh.name.startswith("Content Placeholder") or sh.name == "Rectangle 3":
            return sh
    return None


def object_ph(slide):
    """A non-text OBJECT placeholder -- the template wants a figure there."""
    for sh in slide.shapes:
        if sh.is_placeholder and not sh.has_text_frame                 and sh.name.startswith("Content Placeholder"):
            return sh
    return None


def place_picture(slide, ph, image, max_h_in):
    """Drop `image` into the region `ph` occupies, then remove `ph`."""
    left, top, width = ph.left, ph.top, ph.width
    ph._element.getparent().remove(ph._element)
    pic = slide.shapes.add_picture(str(image), left, top, width=width)
    if pic.height > Inches(max_h_in):
        scale = Inches(max_h_in) / pic.height
        pic.height = Inches(max_h_in)
        pic.width = int(pic.width * scale)
        pic.left = int(left + (width - pic.width) / 2)
    return pic


def main() -> int:
    prs = Presentation(str(SRC))

    # ------------------------------------------------------ title slide ---
    s1 = prs.slides[0]
    for sh in s1.shapes:
        if sh.name == "Title 6":
            tf = sh.text_frame
            set_bullets(tf, [(0, "B.Tech. - BCSE497J Project-I"), (0, TITLE)])
        elif sh.name == "object 3":
            set_bullets(sh.text_frame, [
                (0, "Team member:"),
                (0, "Morisetty Venkata Sai Akhil Teja          (23BCB0135)"),
                (0, ""),
                (0, "Faculty guide:"),
                (0, "Dr. Mythili. T"),
                (0, "School of Computer Science and Engineering"),
            ])

    # ------------------------------------------------- content slides ----
    for num, items in SLIDES.items():
        slide = prs.slides[num - 1]
        ph = content_ph(slide)
        if ph is None:
            print(f"slide {num}: no content placeholder, skipped")
            continue
        set_bullets(ph.text_frame, items)

    # Slide 5's body is an OBJECT placeholder, so the literature survey goes
    # in as the comparison table the template's layout is asking for.
    s5 = prs.slides[4]
    oph = object_ph(s5)
    if oph is not None:
        place_picture(s5, oph, FIG / "fig5_litreview.png", 2.3)

    # SDG / outcomes strip on the Objectives slide
    for sh in prs.slides[6].shapes:
        if sh.name == "TextBox 6":
            set_bullets(sh.text_frame, [(0, SDG_LINE)])

    # --------------------------------- architecture image on slide 8 -----
    s8 = prs.slides[7]
    place_picture(s8, content_ph(s8), FIG / "fig2_architecture.png", 4.7)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print("wrote", OUT)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
